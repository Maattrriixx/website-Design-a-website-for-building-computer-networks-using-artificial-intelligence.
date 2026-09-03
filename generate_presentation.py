from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_THEME_COLOR

OUT = 'NetArchAI_Project_Presentation.pptx'
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(8, 18, 38)
NAVY2 = RGBColor(14, 31, 58)
CYAN = RGBColor(0, 200, 248)
ORANGE = RGBColor(255, 107, 53)
GREEN = RGBColor(16, 185, 129)
WHITE = RGBColor(245, 249, 252)
MUTED = RGBColor(166, 185, 204)
LINE = RGBColor(40, 70, 100)
FONT = 'Arial'


def box(slide, x, y, w, h, fill=NAVY2, line=LINE, radius=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line; shape.line.width = Pt(1)
    return shape


def text(slide, value, x, y, w, h, size=18, color=WHITE, bold=False, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = value; r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb


def title(slide, heading, kicker=None):
    text(slide, heading, 0.65, 0.38, 12.0, 0.58, 27, WHITE, True)
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.95), Inches(0.42), Inches(0.75), Inches(0.08)).fill.solid()
    accent = slide.shapes[-1]; accent.fill.fore_color.rgb = CYAN; accent.line.fill.background()
    if kicker: text(slide, kicker, 0.68, 1.03, 11.7, 0.28, 10, CYAN, True)


def base(heading, kicker=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = NAVY
    for i in range(0, 14):
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(i), Inches(0), Inches(0.008), Inches(7.5))
        ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(12, 29, 52); ln.line.fill.background()
    title(slide, heading, kicker)
    return slide


def bullets(slide, items, x, y, w, h, size=18, color=WHITE, accent=CYAN):
    for i, item in enumerate(items):
        yy = y + i * (h / len(items))
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(yy + 0.09), Inches(0.10), Inches(0.10))
        dot.fill.solid(); dot.fill.fore_color.rgb = accent; dot.line.fill.background()
        text(slide, item, x + 0.22, yy, w - 0.22, (h / len(items)) - 0.05, size, color)


def pill(slide, label, x, y, w, color=CYAN):
    box(slide, x, y, w, 0.38, NAVY2, color, True)
    text(slide, label, x, y + 0.04, w, 0.25, 11, color, True, PP_ALIGN.CENTER)


def node(slide, label, x, y, w, h, color=CYAN, size=14):
    box(slide, x, y, w, h, NAVY2, color, True)
    text(slide, label, x + 0.08, y + 0.07, w - 0.16, h - 0.12, size, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)


def connector(slide, x1, y1, x2, y2, color=LINE):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(1.5)

# 1
s = prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
text(s, 'NetArchAI', 0.75, 1.0, 11.8, 1.0, 48, CYAN, True)
text(s, 'من المخطط المعماري إلى شبكة سلكية ذكية', 0.8, 2.15, 11.4, 0.65, 28, WHITE, True)
text(s, 'عرض مشروع | Laravel + React + AI + Network Optimization', 0.82, 3.0, 11, 0.4, 16, MUTED)
for i, (label, col) in enumerate([('تحليل المخطط', CYAN), ('توليد الشبكة', ORANGE), ('التحرير والتصدير', GREEN)]): pill(s, label, 0.85 + i*2.55, 5.55, 2.15, col)
text(s, 'حل متكامل لتصميم الشبكات اعتمادًا على بيانات المخطط ومساحة الغرف', 0.82, 6.45, 11.5, 0.35, 14, MUTED)

# 2
s = base('مقدمة عن المشروع', '01 | الفكرة والقيمة')
box(s, 0.7, 1.55, 5.65, 4.9, NAVY2, LINE, True)
text(s, 'الفكرة', 1.0, 1.9, 5.0, 0.4, 22, ORANGE, True)
bullets(s, ['رفع مخطط معماري للمبنى', 'اكتشاف الغرف وتصنيفها آليًا', 'اقتراح أجهزة الشبكة وتوزيعها', 'حساب التوصيلات وVLAN تلقائيًا', 'إتاحة التعديل والحفظ والتصدير'], 1.0, 2.55, 4.95, 3.2, 17)
box(s, 6.75, 1.55, 5.85, 4.9, NAVY2, LINE, True)
text(s, 'القيمة المضافة', 7.05, 1.9, 5.2, 0.4, 22, GREEN, True)
text(s, 'تقليل العمل اليدوي وتحويل صورة المخطط إلى تصور شبكة قابل للتنفيذ، مع إبقاء القرار النهائي بيد المصمم.', 7.05, 2.65, 5.1, 1.2, 21, WHITE, True)
pill(s, 'AI-assisted design', 7.05, 4.55, 2.25, CYAN); pill(s, 'Editable topology', 9.55, 4.55, 2.35, ORANGE)

# 3
s = base('المشاكل والصعوبات', '02 | لماذا نحتاج الحل؟')
bullets(s, ['المخطط صورة غير منظمة ولا يحتوي على بيانات شبكة مباشرة.', 'تحديد الغرف ومساحاتها يدويًا بطيء ومعرض للأخطاء.', 'اختيار عدد السويتشات والمنافذ والتوصيلات يحتاج حسابات متكررة.', 'اختلاف أنواع الغرف وصيغ البيانات بين الأنظمة قد يسبب قرارات غير دقيقة.', 'تشغيل Laravel وخدمتي Flask وFastAPI يتطلب تنسيقًا بين بيئات متعددة.'], 0.85, 1.55, 11.7, 4.6, 21, WHITE, ORANGE)
box(s, 0.85, 6.05, 11.7, 0.55, RGBColor(45, 28, 25), ORANGE, True)
text(s, 'المشكلة المركزية: ربط فهم المكان بالتخطيط الشبكي القابل للتنفيذ.', 1.1, 6.18, 11.2, 0.25, 16, ORANGE, True, PP_ALIGN.CENTER)

# 4
s = base('كيف تم حل الصعوبات؟', '02 | منهجية الحل')
steps = [('1', 'اكتشاف', 'YOLO + Keras'), ('2', 'تفسير', 'Rooms + corners'), ('3', 'تحسين', 'OR-Tools CP-SAT'), ('4', 'عرض', 'React Canvas'), ('5', 'تصدير', 'PDF report')]
for i, (num, lab, tech) in enumerate(steps):
    x = 0.8 + i*2.48
    if i < 4: connector(s, x+1.7, 3.0, x+2.35, 3.0, CYAN)
    node(s, num, x+0.62, 1.85, 0.68, 0.68, ORANGE, 22)
    text(s, lab, x, 2.75, 1.95, 0.35, 18, WHITE, True, PP_ALIGN.CENTER)
    text(s, tech, x, 3.25, 1.95, 0.5, 13, MUTED, False, PP_ALIGN.CENTER)
text(s, 'التقسيم إلى خدمات متخصصة جعل دورة العمل قابلة للتوسع والفحص، مع حفظ النتائج في Laravel وإعادة استخدامها في الواجهة.', 1.0, 5.2, 11.2, 0.75, 20, WHITE, True, PP_ALIGN.CENTER)

# 5
s = base('الخدمات المقدمة في المشروع', '03 | ما الذي يقدمه النظام؟')
services = [('إدارة الحساب', 'Register, Login, verification', CYAN), ('إدارة المشاريع', 'إنشاء، عرض، حذف، حالة التحليل', ORANGE), ('تحليل المخطط', 'اكتشاف الغرف والزوايا والمساحات', GREEN), ('تصميم الشبكة', 'أجهزة، VLAN، منافذ، روابط', CYAN), ('لوحة التحكم', 'إحصاءات المستخدمين والمشاريع', ORANGE), ('التصدير', 'تقرير PDF قابل للمشاركة', GREEN)]
for i, (a,b,c) in enumerate(services):
    x = 0.8 + (i%3)*4.15; y = 1.55 + (i//3)*2.2
    box(s, x, y, 3.65, 1.65, NAVY2, c, True); text(s, a, x+0.2, y+0.25, 3.25, 0.35, 20, c, True); text(s, b, x+0.2, y+0.82, 3.25, 0.45, 14, WHITE, False)

# 6
s = base('التقنيات المستخدمة', '04 | Backend, Frontend, Python')
cols = [('Backend', ['Laravel 12 / PHP 8.2', 'Eloquent ORM + migrations', 'Sanctum authentication', 'Guzzle / HTTP Client'], CYAN), ('Frontend', ['React 18 + Vite', 'React Router 6', 'CSS Modules', 'Canvas + jsPDF + AutoTable'], ORANGE), ('Python / AI', ['Flask: hybrid_api.py', 'FastAPI: api_wierd.py', 'YOLO + TensorFlow/Keras', 'OpenCV + NumPy + Shapely', 'OR-Tools CP-SAT'], GREEN)]
for i,(head, items, col) in enumerate(cols):
    x=0.65+i*4.2; box(s,x,1.55,3.75,4.85,NAVY2,col,True); text(s,head,x+0.25,1.88,3.25,0.4,23,col,True,PP_ALIGN.CENTER); bullets(s,items,x+0.35,2.65,3.05,3.15,16,WHITE,col)

# 7
s = base('الميزات التقنية', '04 | منطق التنفيذ')
features = [('اكتشاف مكاني', 'تحليل الصورة واستخراج الغرف والزوايا والمراكز.'), ('تحسين قيودي', 'اختيار مواقع السويتشات وفق المسافة والقيود باستخدام CP-SAT.'), ('شبكة قابلة للتفسير', 'توليد Core Switch وFirewall وRouter والأجهزة الطرفية بروابط واضحة.'), ('بيانات مترابطة', 'حفظ المشروع والغرف والأجهزة والاتصالات في قاعدة بيانات Eloquent.'), ('تجربة تفاعلية', 'Canvas للتكبير والتحريك والتعديل، مع Network/VLAN views.')]
for i,(a,b) in enumerate(features):
    y=1.45+i*0.95; box(s,0.9,y,11.5,0.72,NAVY2,LINE,True); text(s,a,1.2,y+0.13,2.55,0.35,17,CYAN,True); text(s,b,3.95,y+0.13,7.9,0.4,16,WHITE)

# 8
s = base('معمارية النظام', '05 | تدفق البيانات بين الطبقات')
node(s,'React UI',0.75,2.0,2.2,0.72,ORANGE); node(s,'Laravel API',4.0,2.0,2.2,0.72,CYAN); node(s,'MySQL / Eloquent',7.25,2.0,2.45,0.72,GREEN); node(s,'Storage',10.6,2.0,1.75,0.72,ORANGE)
connector(s,2.95,2.36,4.0,2.36,WHITE); connector(s,6.2,2.36,7.25,2.36,WHITE); connector(s,9.7,2.36,10.6,2.36,WHITE)
node(s,'Flask /predict',2.0,4.25,2.35,0.72,CYAN); node(s,'FastAPI /wired',5.45,4.25,2.35,0.72,GREEN); node(s,'AI Models',8.9,4.25,2.35,0.72,ORANGE)
connector(s,5.1,2.72,3.15,4.25,CYAN); connector(s,5.1,2.72,6.62,4.25,GREEN); connector(s,7.8,4.61,8.9,4.61,ORANGE)
text(s,'Laravel هو طبقة التنسيق: يستقبل الطلب، يرسل الصورة أو بيانات الغرف للخدمات المتخصصة، ثم يحفظ النتيجة ويعيدها للواجهة.',1.0,6.05,11.4,0.5,16,WHITE,True,PP_ALIGN.CENTER)

# 9
s = base('رحلة المستخدم داخل النظام', '05 | User flow')
flow = ['تسجيل الدخول', 'إنشاء مشروع', 'رفع المخطط', 'تحليل AI', 'مراجعة الغرف', 'تحسين الشبكة', 'تعديل وحفظ', 'تصدير PDF']
for i,lab in enumerate(flow):
    x=0.55+(i%4)*3.15; y=1.65+(i//4)*2.0
    node(s, f'{i+1}. {lab}', x, y, 2.45, 0.75, [CYAN,ORANGE,GREEN,CYAN][i%4], 15)
    if i not in [3,7]:
        nx= x+2.45 if i%4<3 else x
        ny= y+0.38 if i%4<3 else y+1.25
        connector(s,x+2.45 if i%4<3 else x+1.2,y+0.75 if i%4==3 else y+0.38,nx,ny,MUTED)

# 10
s = base('قاعدة البيانات: Class Diagram', '06 | الجداول والعلاقات الأساسية')
node(s,'User',0.7,1.65,1.8,0.75,CYAN); node(s,'Project',3.25,1.65,1.9,0.75,ORANGE); node(s,'Room',6.0,1.65,1.8,0.75,GREEN); node(s,'Device',8.7,1.65,1.9,0.75,CYAN); node(s,'Connection',11.15,1.65,1.65,0.75,ORANGE)
connector(s,2.5,2.02,3.25,2.02,WHITE); connector(s,5.15,2.02,6.0,2.02,WHITE); connector(s,7.8,2.02,8.7,2.02,WHITE); connector(s,10.6,2.02,11.15,2.02,WHITE)
text(s,'1',2.62,1.7,0.35,0.25,13,MUTED,True,PP_ALIGN.CENTER); text(s,'*',5.45,1.7,0.35,0.25,13,MUTED,True,PP_ALIGN.CENTER); text(s,'*',8.1,1.7,0.35,0.25,13,MUTED,True,PP_ALIGN.CENTER)
node(s,'RoomCorner',4.35,4.15,2.2,0.75,GREEN); connector(s,6.7,2.4,5.45,4.15,GREEN)
text(s,'User 1..* Projects | Project 1..* Rooms / Devices / Connections | Room 1..* RoomCorners وDevices',0.8,5.55,11.8,0.45,17,WHITE,True,PP_ALIGN.CENTER)
text(s,'ملاحظة: connections تحفظ from_device_id وto_device_id لتمثيل طرفي الرابط.',1.0,6.25,11.3,0.3,13,MUTED,False,PP_ALIGN.CENTER)

# 11
s = base('تفصيل الجداول والعلاقات', '06 | نموذج البيانات')
rows = [('users', 'المستخدمون والأدوار والاشتراك', 'hasMany Project'), ('projects', 'بيانات المخطط وحالة التحليل وmetadata', 'belongsTo User; hasMany Room/Device/Connection'), ('rooms', 'المركز، المساحة، النوع، الثقة', 'belongsTo Project; hasMany Corner/Device'), ('room_corners', 'نقاط حدود الغرفة وترتيبها', 'belongsTo Room'), ('devices', 'نوع الجهاز، VLAN، المنافذ، الموقع', 'belongsTo Project وRoom'), ('connections', 'نوع الرابط، السرعة، المسافة، الوسط', 'يربط جهازين')]
for i,(a,b,c) in enumerate(rows):
    y=1.4+i*0.78; box(s,0.7,y,2.25,0.58,NAVY2,CYAN if i%2==0 else ORANGE,True); text(s,a,0.85,y+0.14,1.95,0.25,14,CYAN,True,PP_ALIGN.CENTER); text(s,b,3.25,y+0.12,4.1,0.3,14,WHITE); text(s,c,7.55,y+0.12,5.0,0.3,14,MUTED)

# 12
s = base('واجهات المستخدم', '07 | تجربة المصمم')
node(s,'Sidebar',0.75,1.6,2.25,3.95,ORANGE); text(s,'رفع المخطط\nبيانات المشروع\nإجراءات التحليل\nالأجهزة\nNetwork / VLAN',0.95,2.1,1.85,2.5,17,WHITE,False,PP_ALIGN.CENTER)
box(s,3.35,1.6,7.1,3.95,RGBColor(20,39,62),CYAN,True); text(s,'Canvas Area',3.55,1.82,6.7,0.35,18,CYAN,True,PP_ALIGN.CENTER); box(s,4.0,2.45,5.8,2.25,RGBColor(29,54,77),GREEN,True); text(s,'المخطط + الغرف + الأجهزة + الروابط',4.2,3.27,5.4,0.4,19,WHITE,True,PP_ALIGN.CENTER)
node(s,'Status Bar',10.75,1.6,1.85,3.95,GREEN); text(s,'الغرف\nالأجهزة\nالروابط\nحالة المعالجة',10.98,2.35,1.4,2.2,16,WHITE,False,PP_ALIGN.CENTER)

# 13
s = base('واجهات الـ Dashboard', '07 | المتابعة والإدارة')
for i,(lab,val,col) in enumerate([('المشاريع','Projects',CYAN),('المستخدمون','Users',ORANGE),('الاشتراكات','Subscriptions',GREEN),('النشاط','Activity',CYAN)]):
    x=0.8+i*3.1; box(s,x,1.55,2.65,1.35,NAVY2,col,True); text(s,lab,x+0.15,1.8,2.35,0.28,16,MUTED,True,PP_ALIGN.CENTER); text(s,val,x+0.15,2.23,2.35,0.32,20,col,True,PP_ALIGN.CENTER)
box(s,0.8,3.35,11.9,2.55,NAVY2,LINE,True); text(s,'لوحة تحكم تشغيلية',1.1,3.65,11.2,0.35,21,WHITE,True,PP_ALIGN.CENTER); bullets(s,['عرض إحصاءات النظام والمستخدمين والمشاريع.', 'متابعة حالات المشاريع والتحليل.', 'إدارة المستخدمين والاشتراكات من Laravel API.'],1.45,4.3,10.5,1.25,16,WHITE,GREEN)

# 14
s = base('النتائج والإنجازات', '08 | ما تم تحقيقه')
bullets(s, ['بناء مسار كامل من صورة المخطط إلى تصميم شبكة قابل للتحرير.', 'دمج اكتشاف الغرف مع قياس المساحة والمراكز والزوايا.', 'توليد أجهزة الشبكة والـ VLAN والروابط وفق قواعد وقيود.', 'توفير واجهة تفاعلية للمراجعة والتعديل بدل الاعتماد على ناتج آلي مغلق.', 'حفظ بنية المشروع ونتيجة الشبكة وإتاحة تصدير تقرير PDF.'],0.9,1.55,11.6,4.65,21,WHITE,GREEN)
pill(s,'End-to-end workflow',4.7,6.25,3.0,GREEN)

# 15
s = base('التحديات والحلول المقترحة', '09 | نقاط التحسين')
items=[('تعدد الخدمات', 'توثيق أوامر التشغيل وإدارة متغيرات البيئة والمنافذ.'),('اختلاف صيغ البيانات', 'توحيد أسماء الحقول وأنواع الغرف بعقد API واضح.'),('مسارات قديمة', 'مراجعة NetworkController وAiAPI مقابل المسارات الحالية.'),('اتساق الصور', 'توحيد image/thumbnail مع Accessors وStorage URLs.'),('الاختبار', 'إضافة اختبارات API وتكامل لخدمة التحليل والتحسين.')]
for i,(a,b) in enumerate(items):
    y=1.4+i*0.92; text(s,a,1.0,y,2.55,0.35,18,ORANGE,True); text(s,b,3.85,y,8.2,0.4,17,WHITE); slide_line=slide_line if False else None
    ln=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(y+0.68), Inches(11.35), Inches(0.015)); ln.fill.solid(); ln.fill.fore_color.rgb=LINE; ln.line.fill.background()

# 16
s = base('ملخص سريع عن المشروع', '10 | في دقيقة واحدة')
text(s,'NetArchAI هو نظام يساعد مصمم الشبكات على تحويل المخطط المعماري إلى شبكة سلكية منظمة.',1.0,1.55,11.3,0.65,25,WHITE,True,PP_ALIGN.CENTER)
summary=[('يدخل', 'مخطط معماري + بيانات المشروع'),('يعالج', 'AI لاكتشاف الغرف وOptimization لتوزيع الشبكة'),('ينتج', 'أجهزة، VLAN، روابط، metadata وتقرير PDF')]
for i,(a,b) in enumerate(summary):
    x=0.95+i*4.1; box(s,x,3.0,3.55,1.6,NAVY2,[CYAN,ORANGE,GREEN][i],True); text(s,a,x+0.2,3.3,3.15,0.3,19,[CYAN,ORANGE,GREEN][i],True,PP_ALIGN.CENTER); text(s,b,x+0.25,3.85,3.05,0.42,15,WHITE,False,PP_ALIGN.CENTER)
text(s,'النتيجة: قرار أسرع، بيانات أوضح، وتصميم يمكن مراجعته وتعديله.',1.1,5.65,11.1,0.4,20,GREEN,True,PP_ALIGN.CENTER)

# 17
s = base('الخلاصة التقنية', '10 | لماذا التصميم قابل للتوسع؟')
bullets(s,['فصل الواجهة عن API وعن خدمات الذكاء الاصطناعي.', 'قاعدة بيانات تحفظ الكيانات الأساسية والعلاقات التشغيلية.', 'خدمة التحسين قابلة لتغيير القيود وقواعد توزيع الأجهزة.', 'واجهة Canvas تسمح بالمراجعة البشرية بعد الاقتراح الآلي.', 'البنية الحالية تفتح الطريق لإضافة تقارير، أنواع شبكات، ومصادر AI جديدة.'],1.0,1.65,11.3,4.5,21,WHITE,CYAN)

# 18
s = prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
text(s,'شكرًا',0.8,2.0,11.7,0.9,42,CYAN,True,PP_ALIGN.CENTER)
text(s,'NetArchAI | Intelligent wired network design',1.0,3.15,11.3,0.45,20,WHITE,True,PP_ALIGN.CENTER)
text(s,'Laravel  •  React  •  Python AI  •  Network Optimization',1.0,4.0,11.3,0.35,15,MUTED,False,PP_ALIGN.CENTER)

for idx, slide in enumerate(prs.slides, 1):
    footer = slide.shapes.add_textbox(Inches(0.65), Inches(7.12), Inches(12), Inches(0.18))
    p=footer.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
    r=p.add_run(); r.text=f'NETARCHAI  |  {idx:02d}'; r.font.name=FONT; r.font.size=Pt(8); r.font.color.rgb=MUTED

prs.save(OUT)
print(f'created {OUT} with {len(prs.slides)} slides')
